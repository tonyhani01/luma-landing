#!/usr/bin/env python3
"""Faithful-ish PIL rasterizer for the generated pptx (preview/QA only)."""
import io, re
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DPI = 150
EMU = 914400
def px(emu): return int(round(emu/EMU*DPI))
def ptpx(pt): return max(1,int(round(pt*DPI/72.0)))

FONTS = {
 'Sora':            '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
 'Sora|b':          '/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf',
 'Open Sans':       '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
 'Open Sans|b':     '/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf',
 'Instrument Serif':'/usr/share/fonts/truetype/liberation/LiberationSerif-Italic.ttf',
 'Instrument Serif|i':'/usr/share/fonts/truetype/liberation/LiberationSerif-Italic.ttf',
 'JetBrains Mono':  '/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf',
 'JetBrains Mono|b':'/usr/share/fonts/truetype/dejavu/DejaVuSansMono-Bold.ttf',
}
_cache={}
def font(name,size,bold,italic):
    key=name+('|b' if bold else '')+('|i' if italic else '')
    path=FONTS.get(key) or FONTS.get(name) or FONTS['Open Sans']
    ck=(path,size)
    if ck not in _cache: _cache[ck]=ImageFont.truetype(path,size)
    return _cache[ck]

# strip emoji / symbols that Liberation lacks (preview clarity); keep latin+arabic+punct
EMOJI = re.compile("[\U0001F000-\U0001FAFF\U00002600-\U000027BF\U00002190-\U000021FF\U00002B00-\U00002BFF\U0000FE00-\U0000FE0F•]")
def clean(t):
    t=EMOJI.sub('',t)
    return t.replace('  ',' ').strip()

def alpha_of(clr_el):
    if clr_el is None: return 255
    a=clr_el.find(qn('a:alpha'))
    if a is not None: return int(int(a.get('val'))/1000/100*255)
    return 255

def shape_fill(sp):
    f=sp.fill
    try:
        if f.type is None: return None
    except Exception: return None
    el=sp._element.spPr.find(qn('a:solidFill'))
    if el is None: return None
    srgb=el.find(qn('a:srgbClr'))
    if srgb is None: return None
    rgb=tuple(int(srgb.get('val')[i:i+2],16) for i in (0,2,4))
    return rgb+(alpha_of(srgb),)

def shape_line(sp):
    ln=sp._element.spPr.find(qn('a:ln'))
    if ln is None: return None,0
    sf=ln.find(qn('a:solidFill'))
    if sf is None: return None,0
    srgb=sf.find(qn('a:srgbClr'))
    if srgb is None: return None,0
    rgb=tuple(int(srgb.get('val')[i:i+2],16) for i in (0,2,4))
    w=ln.get('w'); wpt=int(w)/12700 if w else 1
    return rgb+(alpha_of(srgb),), max(1,ptpx(wpt))

def geom(sp):
    g=sp._element.spPr.find(qn('a:prstGeom'))
    return g.get('prst') if g is not None else 'rect'

def run_props(r):
    rPr=r._r.find(qn('a:rPr'))
    sz=18.0; bold=False; ital=False; name='Open Sans'; col=(255,255,255)
    if rPr is not None:
        if rPr.get('sz'): sz=int(rPr.get('sz'))/100
        bold=rPr.get('b')=='1'; ital=rPr.get('i')=='1'
        lt=rPr.find(qn('a:latin'))
        if lt is not None and lt.get('typeface'): name=lt.get('typeface')
        sf=rPr.find(qn('a:solidFill'))
        if sf is not None:
            srgb=sf.find(qn('a:srgbClr'))
            if srgb is not None: col=tuple(int(srgb.get('val')[i:i+2],16) for i in (0,2,4))
    return sz,bold,ital,name,col

def wrap_runs(runs, fnt_runs, maxw, draw):
    """runs:list of (text,font,color). Greedy wrap into lines preserving runs."""
    words=[]
    for (txt,fnt,col) in fnt_runs:
        parts=re.split('(\\s+)',txt)
        for pa in parts:
            if pa=='':continue
            words.append((pa,fnt,col))
    lines=[[]]; widths=[0]
    for w,fnt,col in words:
        ww=draw.textlength(w,font=fnt)
        if widths[-1]+ww>maxw and w.strip() and lines[-1]:
            lines.append([]); widths.append(0)
        if w.strip()=='' and not lines[-1]:
            continue
        lines[-1].append((w,fnt,col)); widths[-1]+=ww
    return lines

def render_slide(slide, idx):
    img=Image.new('RGB',(px(SW),px(SH)),(5,2,26))
    draw=ImageDraw.Draw(img,'RGBA')
    for sp in slide.shapes:
        try: L,T,W,H=px(sp.left),px(sp.top),px(sp.width),px(sp.height)
        except: continue
        st=sp.shape_type
        if st==13:  # picture
            blob=sp.image.blob
            pim=Image.open(io.BytesIO(blob)).convert('RGBA')
            pim=pim.resize((max(1,W),max(1,H)))
            img.paste(pim,(L,T),pim)
            continue
        # autoshape fill/line
        g=geom(sp)
        fill=shape_fill(sp); line,lw=shape_line(sp)
        if g in ('ellipse',):
            if fill: draw.ellipse([L,T,L+W,T+H],fill=fill)
            if line: draw.ellipse([L,T,L+W,T+H],outline=line,width=lw)
        elif g in ('roundRect',):
            rad=int(min(W,H)*0.5)
            try: rad=int(min(W,H)*float(sp.adjustments[0]))
            except: pass
            rad=max(2,min(rad,int(min(W,H)/2)))
            if fill: draw.rounded_rectangle([L,T,L+W,T+H],radius=rad,fill=fill)
            if line: draw.rounded_rectangle([L,T,L+W,T+H],radius=rad,outline=line,width=lw)
        else:
            if fill: draw.rectangle([L,T,L+W,T+H],fill=fill)
            if line: draw.rectangle([L,T,L+W,T+H],outline=line,width=lw)
        # text
        if sp.has_text_frame and sp.text_frame.text.strip():
            tf=sp.text_frame
            ml=px(tf.margin_left or 0); mr=px(tf.margin_right or 0)
            mt=px(tf.margin_top or 0); mb=px(tf.margin_bottom or 0)
            tx=L+ml; tw=W-ml-mr
            paras=[]
            for p in tf.paragraphs:
                align=p.alignment
                frs=[]
                for r in p.runs:
                    if not r.text: continue
                    sz,b,i,nm,col=run_props(r)
                    frs.append((clean(r.text),font(nm,ptpx(sz),b,i),col,ptpx(sz)))
                paras.append((align,frs,p.line_spacing))
            # measure total height
            blocks=[]
            for align,frs,ls in paras:
                if not frs:
                    blocks.append((align,[ ],ptpx(10))); continue
                fr2=[(t,f,c) for (t,f,c,s) in frs]
                maxsz=max(s for (_,_,_,s) in frs)
                lines=wrap_runs(None,fr2,tw,draw)
                lh=int(maxsz*1.18*(ls if ls else 1))
                blocks.append((align,lines,lh))
            total=sum(len(l) if l else 1 for (a,l,h) in blocks for _ in [0])  # placeholder
            total_h=sum((len(l)*h if l else h) for (a,l,h) in blocks)
            anchor=tf.vertical_anchor
            if anchor==MSO_ANCHOR.MIDDLE: cy=T+mt+ max(0,(H-mt-mb-total_h)//2)
            elif anchor==MSO_ANCHOR.BOTTOM: cy=T+H-mb-total_h
            else: cy=T+mt
            for align,lines,lh in blocks:
                if not lines:
                    cy+=lh; continue
                for line in lines:
                    lw_tot=sum(draw.textlength(t,font=f) for (t,f,c) in line)
                    if align==PP_ALIGN.CENTER: xstart=tx+(tw-lw_tot)/2
                    elif align==PP_ALIGN.RIGHT: xstart=tx+tw-lw_tot
                    else: xstart=tx
                    xx=xstart
                    for (t,f,c) in line:
                        draw.text((xx,cy),t,font=f,fill=c)
                        xx+=draw.textlength(t,font=f)
                    cy+=lh
    return img

prs=Presentation('proposal/LUMA_Business_Proposal.pptx')
SW,SH=prs.slide_width,prs.slide_height
imgs=[render_slide(s,i) for i,s in enumerate(prs.slides)]
# contact sheet: 3 cols x 4 rows
cols,rows=3,4; pad=24
tw,th=imgs[0].size
sheet=Image.new('RGB',(cols*tw+(cols+1)*pad, rows*th+(rows+1)*pad),(15,15,25))
for i,im in enumerate(imgs):
    c=i%cols; r=i//cols
    sheet.paste(im,(pad+c*(tw+pad), pad+r*(th+pad)))
sheet.save('proposal/preview_contact_sheet.png')
# also save individual full-res for slides 1,3,5,10
for n in (0,2,4,9,11):
    imgs[n].save(f'proposal/preview_slide_{n+1:02d}.png')
print("rendered",len(imgs),"slides -> contact sheet + individual previews")
