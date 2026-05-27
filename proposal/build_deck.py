#!/usr/bin/env python3
"""Generates the LUMA business proposal as an interactive .pptx deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
from lxml import etree

# ---------- Brand ----------
C = {
    'cosmos':  RGBColor(0x02,0x01,0x08),
    'cosmos2': RGBColor(0x08,0x03,0x30),
    'glow':    RGBColor(0x44,0x66,0xff),
    'glow_l':  RGBColor(0x7a,0xa8,0xff),
    'glow_xl': RGBColor(0xb3,0xcc,0xff),
    'white':   RGBColor(0xff,0xff,0xff),
    'ink':     RGBColor(0x0a,0x0a,0x2e),
    'mute_d':  RGBColor(0xb3,0xc2,0xe0),  # muted on dark
    'mute_l':  RGBColor(0x55,0x60,0x80),  # muted on light
    'surface': RGBColor(0xf8,0xfa,0xff),
    'green':   RGBColor(0x22,0xc5,0x5e),
    'cardln':  RGBColor(0x44,0x66,0xff),
}
F_HEAD = 'Sora'
F_BODY = 'Open Sans'
F_SERIF= 'Instrument Serif'
F_MONO = 'JetBrains Mono'

ASSETS = 'proposal/assets'
TOTAL = 14

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

# ---------- helpers ----------
def slide(bg):
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(f'{ASSETS}/{bg}', 0, 0, SW, SH)
    return s

def _set_font(run, size, color, font=F_BODY, bold=False, italic=False, spacing=None):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic
    if spacing is not None:
        run.font._rPr.set('spc', str(int(spacing*100)))

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=None, space_after=None, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt, kwargs)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        if space_after is not None: p.space_after = Pt(space_after)
        for txt, kw in para:
            r = p.add_run(); r.text = txt
            _set_font(r, **kw)
    return tb

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.10, shadow=False):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        _soft_shadow(sp)
    return sp

def set_alpha(shape, alpha_pct):
    """Set fill transparency (0-100)."""
    sf = shape.fill.fore_color._xFill
    srgb = sf.find(qn('a:srgbClr'))
    a = srgb.makeelement(qn('a:alpha'), {'val': str(int((100-alpha_pct)*1000))})
    srgb.append(a)

def _soft_shadow(sp):
    spPr = sp._element.spPr
    el = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(el, qn('a:outerShdw'),
        {'blurRad':'180000','dist':'90000','dir':'5400000','rotWithShape':'0'})
    c = etree.SubElement(sh, qn('a:srgbClr'), {'val':'1A1A6E'})
    etree.SubElement(c, qn('a:alpha'), {'val':'34000'})

def shape_text(sp, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for txt, kw in para:
            r = p.add_run(); r.text = txt
            _set_font(r, **kw)

def eyebrow(s, x, y, label, color=None):
    color = color or C['glow_l']
    text(s, x, y, Inches(8), Inches(0.4),
         [[(label.upper(), dict(size=12.5, color=color, font=F_MONO, bold=True, spacing=2.6))]])

# ---------- transitions (varied, non-repetitive) ----------
def transition(s, kind, **attrs):
    sld = s._element
    # remove existing
    for ex in sld.findall(qn('p:transition')):
        sld.remove(ex)
    t = etree.SubElement(sld, qn('p:transition'))
    t.set('spd', attrs.pop('spd', 'med'))
    child = etree.SubElement(t, qn('p:'+kind))
    for k, v in attrs.items():
        child.set(k, v)
    # position: must come after clrMapOvr
    cmo = sld.find(qn('p:clrMapOvr'))
    if cmo is not None:
        cmo.addnext(t)

# ---------- internal hyperlinks (interactivity) ----------
_pending_links = []
def link_to_slide(shape, target_index):
    """Defer until all slides exist (handles forward references)."""
    _pending_links.append((shape, target_index))

def _resolve_links():
    for shape, target_index in _pending_links:
        _do_link_to_slide(shape, target_index)

def _do_link_to_slide(shape, target_index):
    """Make a shape jump to slide at target_index (0-based) on click."""
    target = prs.slides[target_index]
    part = shape.part
    rId = part.relate_to(target.part,
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
    cNvPr = shape._element._nvXxPr.cNvPr
    for ex in cNvPr.findall(qn('a:hlinkClick')):
        cNvPr.remove(ex)
    hlink = etree.SubElement(cNvPr, qn('a:hlinkClick'))
    hlink.set(qn('r:id'), rId)
    hlink.set('action', 'ppaction://hlinksldjump')

def link_url(shape, url):
    cNvPr = shape._element._nvXxPr.cNvPr
    part = shape.part
    rId = part.relate_to(url,
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
        is_external=True)
    h = etree.SubElement(cNvPr, qn('a:hlinkClick'))
    h.set(qn('r:id'), rId)

def home_button(s, target_index=1):
    """Small 'menu' chip bottom-left that jumps to the contents slide."""
    b = rect(s, Inches(0.45), Inches(6.92), Inches(1.15), Inches(0.42),
             fill=None, line=C['glow_l'], line_w=1.0, radius=0.5)
    set_alpha_line(b, 30)
    shape_text(b, [[('‹  Menu', dict(size=10.5, color=C['glow_l'], font=F_MONO, bold=True))]],
               align=PP_ALIGN.CENTER)
    link_to_slide(b, target_index)

def set_alpha_line(shape, alpha_pct):
    ln = shape.line._get_or_add_ln()
    sf = ln.find(qn('a:solidFill'))
    srgb = sf.find(qn('a:srgbClr'))
    etree.SubElement(srgb, qn('a:alpha'), {'val': str(int((100-alpha_pct)*1000))})

def page_num(s, n, dark=True):
    col = C['mute_d'] if dark else C['mute_l']
    text(s, Inches(11.9), Inches(6.95), Inches(1.0), Inches(0.4),
         [[(f'{n:02d} / {TOTAL}', dict(size=10, color=col, font=F_MONO))]], align=PP_ALIGN.RIGHT)

def brand_mark(s, dark=True):
    s.shapes.add_picture(f'{ASSETS}/../../logo.png', Inches(0.5), Inches(0.42),
                         height=Inches(0.5))

# =====================================================================
# SLIDE 1 — COVER
# =====================================================================
s = slide('bg_cover.png')
s.shapes.add_picture('logo.png', Inches(5.62), Inches(1.15), height=Inches(1.45))
text(s, Inches(1), Inches(2.62), Inches(11.33), Inches(0.5),
     [[('LUMA', dict(size=30, color=C['white'], font=F_HEAD, bold=True, spacing=8))]],
     align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(3.18), Inches(11.33), Inches(1.4),
     [[('The AI teammate behind ', dict(size=40, color=C['white'], font=F_HEAD, bold=True)),
       ('every DM', dict(size=40, color=C['glow_l'], font=F_HEAD, bold=True))],
      [('— in the language they wrote to you in.', dict(size=27, color=C['glow_xl'], font=F_SERIF, italic=True))]],
     align=PP_ALIGN.CENTER, line_spacing=1.1, space_after=6)
# pill
pill = rect(s, Inches(3.77), Inches(4.95), Inches(5.8), Inches(0.5), fill=C['cosmos2'],
            line=C['glow'], line_w=1.0, radius=0.5)
set_alpha(pill, 45); set_alpha_line(pill, 35)
shape_text(pill, [[('BUSINESS PROPOSAL  ·  AI CUSTOMER SERVICE', dict(size=12.5, color=C['glow_xl'], font=F_MONO, bold=True, spacing=1.6))]], align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(5.95), Inches(11.33), Inches(0.8),
     [[('Prepared for forward-thinking local businesses', dict(size=15, color=C['mute_d'], font=F_BODY))],
      [('luma-bot.com', dict(size=13, color=C['glow_l'], font=F_MONO))]],
     align=PP_ALIGN.CENTER, space_after=4)
# clickable explore button
cta = rect(s, Inches(5.46), Inches(6.5), Inches(2.4), Inches(0.52), fill=C['glow'], radius=0.5, shadow=True)
shape_text(cta, [[('Explore the proposal  →', dict(size=13, color=C['white'], font=F_HEAD, bold=True))]], align=PP_ALIGN.CENTER)
link_to_slide(cta, 1)
transition(s, 'fade', spd='slow')

# =====================================================================
# SLIDE 2 — CONTENTS (interactive)
# =====================================================================
s = slide('bg_dark.png')
brand_mark(s)
eyebrow(s, Inches(0.85), Inches(1.15), 'Navigation')
text(s, Inches(0.8), Inches(1.5), Inches(8), Inches(1),
     [[('What’s inside', dict(size=44, color=C['white'], font=F_HEAD, bold=True))]])
toc = [
    ('01', 'The problem', 'DMs pile up, replies lag, sales slip away', 2),
    ('02', 'Meet LUMA', 'Your always-on AI teammate', 3),
    ('03', 'How it works', 'Answer · Reply · Hand off', 4),
    ('04', 'Get started in 3 steps', 'Connect, teach, go live', 5),
    ('05', 'Feature set', 'Everything in one inbox', 6),
    ('06', 'Insights & analytics', 'Know your customers at a glance', 7),
    ('07', 'Why LUMA', 'Built for merchants, not generic bots', 9),
    ('08', 'Proven performance', 'Speed, languages, uptime', 10),
    ('09', 'Pricing', 'Plans that scale with you', 11),
    ('10', 'Security & trust', 'Your data, protected', 12),
    ('11', 'Next steps', 'Start your free trial', 13),
]
col_x = [Inches(0.8), Inches(6.95)]
y0 = Inches(2.42); rh = Inches(0.78)
for i, (num, title, sub, tgt) in enumerate(toc):
    col = i // 6; row = i % 6
    x = col_x[col]; y = Emu(int(y0) + row*int(rh))
    card = rect(s, x, y, Inches(5.55), Inches(0.66), fill=C['cosmos2'], line=C['glow'], line_w=0.75, radius=0.18)
    set_alpha(card, 55); set_alpha_line(card, 55)
    link_to_slide(card, tgt)
    text(s, Emu(int(x)+Emu(Inches(0.22))), Emu(int(y)+Emu(Inches(0.12))), Inches(0.9), Inches(0.5),
         [[(num, dict(size=22, color=C['glow_l'], font=F_MONO, bold=True))]])
    text(s, Emu(int(x)+Emu(Inches(1.15))), Emu(int(y)+Emu(Inches(0.09))), Inches(4.2), Inches(0.6),
         [[(title, dict(size=15.5, color=C['white'], font=F_HEAD, bold=True))],
          [(sub, dict(size=10.5, color=C['mute_d'], font=F_BODY))]], space_after=1)
text(s, Inches(0.8), Inches(7.0), Inches(8), Inches(0.4),
     [[('Tip: every card and button in this deck is clickable.', dict(size=10.5, color=C['glow_l'], font=F_MONO))]])
page_num(s, 2)
transition(s, 'push', dir='u')

# =====================================================================
# SLIDE 3 — THE PROBLEM
# =====================================================================
s = slide('bg_dark2.png')
eyebrow(s, Inches(0.85), Inches(0.85), '01 · The problem')
text(s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.6),
     [[('Every silent DM is a ', dict(size=42, color=C['white'], font=F_HEAD, bold=True)),
       ('sale walking out the door', dict(size=42, color=C['glow_l'], font=F_HEAD, bold=True)),
       ('.', dict(size=42, color=C['white'], font=F_HEAD, bold=True))]], line_spacing=1.05)
probs = [
    ('Replies come too late', 'Customers message at midnight and on weekends. By the time you reply, they’ve bought elsewhere.'),
    ('The same questions, endlessly', '“Is this in stock?” “Where’s my order?” “Do you ship to me?” — answered by hand, all day long.'),
    ('Customers write in three languages', 'Arabic, English and Franco-Arabic in the same inbox. Your team juggles tone and translation.'),
    ('No view of what’s working', 'No idea how many messages turn into sales, what gets asked most, or when your busiest hours are.'),
]
gx = Inches(0.8); gy = Inches(2.95); cw = Inches(5.85); ch = Inches(1.78); gap = Inches(0.3)
for i,(t,d) in enumerate(probs):
    col=i%2; row=i//2
    x=Emu(int(gx)+col*(int(cw)+int(gap))); y=Emu(int(gy)+row*(int(ch)+int(gap)))
    card=rect(s,x,y,cw,ch,fill=C['cosmos2'],line=C['glow'],line_w=0.75,radius=0.10)
    set_alpha(card,50); set_alpha_line(card,60)
    text(s,Emu(int(x)+Emu(Inches(0.32))),Emu(int(y)+Emu(Inches(0.22))),Inches(1.1),Inches(0.7),
         [[(f'{i+1:02d}',dict(size=30,color=C['glow_l'],font=F_MONO,bold=True))]])
    text(s,Emu(int(x)+Emu(Inches(1.3))),Emu(int(y)+Emu(Inches(0.28))),Emu(int(cw)-Emu(Inches(1.6))),Inches(1.2),
         [[(t,dict(size=17,color=C['white'],font=F_HEAD,bold=True))],
          [(d,dict(size=12,color=C['mute_d'],font=F_BODY))]],line_spacing=1.05,space_after=4)
home_button(s); page_num(s,3)
transition(s, 'wipe', dir='l')

# =====================================================================
# SLIDE 4 — MEET LUMA
# =====================================================================
s = slide('bg_dark.png')
eyebrow(s, Inches(0.85), Inches(0.85), '02 · Meet LUMA')
text(s, Inches(0.8), Inches(1.25), Inches(7.4), Inches(2.6),
     [[('An AI teammate that ', dict(size=40, color=C['white'], font=F_HEAD, bold=True)),
       ('never sleeps', dict(size=40, color=C['glow_l'], font=F_HEAD, bold=True))],
      [('', dict(size=8, color=C['white'], font=F_BODY))],
      [('LUMA reads every Instagram message and website chat, answers in the customer’s own language using your live product and order information, and passes the conversation to your team whenever a personal touch is needed.',
        dict(size=16.5, color=C['mute_d'], font=F_BODY))]], line_spacing=1.12, space_after=6)
bullets = [
    'Replies instantly, day and night',
    'Speaks Egyptian Arabic, English and Franco-Arabic',
    'Answers order and product questions accurately',
    'Works inside your store — no new app to learn',
]
by = Inches(3.95)
for b in bullets:
    dot = rect(s, Inches(0.85), Emu(int(by)+Emu(Inches(0.06))), Inches(0.16), Inches(0.16), fill=C['glow'], shape=MSO_SHAPE.OVAL)
    text(s, Inches(1.2), by, Inches(6.6), Inches(0.5),
         [[(b, dict(size=14.5, color=C['white'], font=F_BODY))]])
    by = Emu(int(by)+Emu(Inches(0.62)))
# chat mock card
mx=Inches(8.55); mw=Inches(4.1)
mock=rect(s,mx,Inches(1.25),mw,Inches(5.1),fill=C['cosmos2'],line=C['glow'],line_w=1.0,radius=0.06,shadow=True)
set_alpha(mock,18); set_alpha_line(mock,55)
hdr=rect(s,mx,Inches(1.25),mw,Inches(0.78),fill=C['glow'],radius=0.06)
set_alpha(hdr,15)
s.shapes.add_picture('logo.png', Emu(int(mx)+Emu(Inches(0.22))), Inches(1.43), height=Inches(0.42))
text(s,Emu(int(mx)+Emu(Inches(0.85))),Inches(1.4),Inches(3),Inches(0.4),
     [[('LUMA',dict(size=13,color=C['white'],font=F_HEAD,bold=True))]])
odot=rect(s,Emu(int(mx)+Emu(Inches(0.86))),Inches(1.79),Inches(0.085),Inches(0.085),fill=C['green'],shape=MSO_SHAPE.OVAL)
text(s,Emu(int(mx)+Emu(Inches(1.02))),Inches(1.71),Inches(2),Inches(0.3),
     [[('Online now',dict(size=9,color=C['green'],font=F_MONO))]])
chat=[('in','هل في الجاكيت الأسود مقاس M؟',C['cosmos']),
      ('out','Yes! The black jacket is in stock in size M — EGP 1,250. Would you like me to reserve one?',C['glow']),
      ('in','Where’s my order #4521?',C['cosmos']),
      ('out','It’s out for delivery today and should arrive by 6 PM.',C['glow'])]
cy=Inches(2.25)
for side,msg,col in chat:
    bw=Inches(3.0)
    bx = Emu(int(mx)+Emu(Inches(0.85))) if side=='in' else Emu(int(mx)+Emu(Inches(0.25)))
    if side=='out': bx=Emu(int(mx)+int(mw)-int(bw)-Emu(Inches(0.25)))
    bub=rect(s,bx,cy,bw,Inches(0.92),fill=col,radius=0.22)
    if side=='in': set_alpha(bub,25)
    shape_text(bub,[[(msg,dict(size=10.5,color=C['white'],font=F_BODY))]],
               align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.MIDDLE)
    cy=Emu(int(cy)+Emu(Inches(1.02)))
home_button(s); page_num(s,4)
transition(s, 'cover', dir='r')

# =====================================================================
# SLIDE 5 — HOW IT WORKS (3 outcomes)
# =====================================================================
s = slide('bg_dark2.png')
eyebrow(s, Inches(0.85), Inches(0.7), '03 · How it works')
text(s, Inches(0.8), Inches(1.05), Inches(11.5), Inches(0.9),
     [[('Three ways LUMA handles every message', dict(size=34, color=C['white'], font=F_HEAD, bold=True))]])
text(s, Inches(0.8), Inches(1.75), Inches(11.5), Inches(0.5),
     [[('Every message is handled the right way — keeping replies instant and your customers looked after.', dict(size=14, color=C['mute_d'], font=F_BODY))]])
cards=[('Instant answer','Greetings, thank-yous and repeat questions are recognised and answered straight away.',
        ['“Hi there”','“شكراً”','“Where’s order #4521?”']),
       ('Smart reply','Real product questions are answered from your live catalogue, in the customer’s own language.',
        ['“Is the black jacket in M?”','“Do you ship to Alexandria?”','“What’s the return policy?”']),
       ('Personal handoff','Complaints or anything unusual? LUMA reassures the customer and passes it straight to your team.',
        ['“I’d like to speak to someone”','“This ripped after one wear”','“I’m not happy”'])]
cw=Inches(3.95); gap=Inches(0.27); x0=Inches(0.8); y=Inches(2.5); ch=Inches(4.05)
for i,(t,d,ex) in enumerate(cards):
    x=Emu(int(x0)+i*(int(cw)+int(gap)))
    card=rect(s,x,y,cw,ch,fill=C['cosmos2'],line=C['glow'],line_w=0.9,radius=0.07,shadow=True)
    set_alpha(card,45); set_alpha_line(card,50)
    # number badge
    bdg=rect(s,Emu(int(x)+Emu(Inches(0.3))),Emu(int(y)+Emu(Inches(0.3))),Inches(0.85),Inches(0.85),fill=C['glow'],shape=MSO_SHAPE.OVAL)
    shape_text(bdg,[[(f'{i+1}',dict(size=30,color=C['white'],font=F_HEAD,bold=True))]],align=PP_ALIGN.CENTER)
    text(s,Emu(int(x)+Emu(Inches(0.3))),Emu(int(y)+Emu(Inches(1.4))),Emu(int(cw)-Emu(Inches(0.6))),Inches(1.6),
         [[(t,dict(size=20,color=C['white'],font=F_HEAD,bold=True))],
          [(d,dict(size=12,color=C['mute_d'],font=F_BODY))]],line_spacing=1.08,space_after=4)
    ey=Emu(int(y)+Emu(Inches(2.95)))
    for e in ex:
        chip=rect(s,Emu(int(x)+Emu(Inches(0.3))),ey,Emu(int(cw)-Emu(Inches(0.6))),Inches(0.32),fill=C['cosmos'],line=C['glow'],line_w=0.5,radius=0.5)
        set_alpha(chip,40); set_alpha_line(chip,70)
        shape_text(chip,[[(e,dict(size=9.5,color=C['glow_xl'],font=F_MONO))]],align=PP_ALIGN.LEFT)
        ey=Emu(int(ey)+Emu(Inches(0.37)))
home_button(s); page_num(s,5)
transition(s, 'split', orient='vert', dir='out')

# =====================================================================
# SLIDE 6 — GET STARTED IN 3 STEPS
# =====================================================================
s = slide('bg_light.png')
eyebrow(s, Inches(0.85), Inches(0.85), '04 · Onboarding', color=C['glow'])
text(s, Inches(0.8), Inches(1.25), Inches(11.5), Inches(1),
     [[('Live in an afternoon — ', dict(size=38, color=C['ink'], font=F_HEAD, bold=True)),
       ('no setup headaches', dict(size=38, color=C['glow'], font=F_HEAD, bold=True))]])
steps=[('1','Connect','Link your store and Instagram account in a few clicks. Your products load in automatically.'),
       ('2','Personalise','Set the tone, your opening hours and a few product details so LUMA sounds just like your brand.'),
       ('3','Go live','Watch conversations as they happen and step into any chat yourself whenever you’d like.')]
cw=Inches(3.85); gap=Inches(0.4); x0=Inches(0.8); y=Inches(3.0); ch=Inches(3.1)
for i,(n,t,d) in enumerate(steps):
    x=Emu(int(x0)+i*(int(cw)+int(gap)))
    card=rect(s,x,y,cw,ch,fill=C['white'],line=C['glow'],line_w=1.0,radius=0.08,shadow=True)
    set_alpha_line(card,80)
    circ=rect(s,Emu(int(x)+Emu(Inches(0.35))),Emu(int(y)+Emu(Inches(0.4))),Inches(1.0),Inches(1.0),fill=C['glow'],shape=MSO_SHAPE.OVAL,shadow=True)
    shape_text(circ,[[(n,dict(size=34,color=C['white'],font=F_HEAD,bold=True))]],align=PP_ALIGN.CENTER)
    text(s,Emu(int(x)+Emu(Inches(0.35))),Emu(int(y)+Emu(Inches(1.65))),Emu(int(cw)-Emu(Inches(0.7))),Inches(1.4),
         [[(t,dict(size=22,color=C['ink'],font=F_HEAD,bold=True))],
          [(d,dict(size=12.5,color=C['mute_l'],font=F_BODY))]],line_spacing=1.1,space_after=6)
    if i<2:
        text(s,Emu(int(x)+int(cw)-Emu(Inches(0.05))),Emu(int(y)+Emu(Inches(1.15))),Inches(0.5),Inches(0.6),
             [[('→',dict(size=30,color=C['glow_l'],font=F_HEAD,bold=True))]],align=PP_ALIGN.CENTER)
home_button(s); page_num(s,6,dark=False)
transition(s, 'push', dir='l')

# =====================================================================
# SLIDE 7 — FEATURE SET
# =====================================================================
s = slide('bg_dark.png')
eyebrow(s, Inches(0.85), Inches(0.6), '05 · Feature set')
text(s, Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.8),
     [[('Everything your inbox needs, in one place', dict(size=33, color=C['white'], font=F_HEAD, bold=True))]])
feats=[('One inbox','Instagram and website chats together in one place'),
       ('Three languages','Speaks Arabic, English and Franco-Arabic'),
       ('Order tracking','Tells customers exactly where their order is'),
       ('Instant answers','Common questions answered in a moment'),
       ('Smart handoff','Knows when to bring in your team'),
       ('Simple dashboard','See your busiest hours and top questions'),
       ('Product knowledge','Shares accurate details about your products'),
       ('After-hours replies','Keeps answering when you’re closed'),
       ('Customer feedback','Asks each customer if the answer helped')]
cw=Inches(3.85); ch=Inches(1.55); gx=Inches(0.32); gy=Inches(0.3); x0=Inches(0.8); y0=Inches(1.95)
for i,(t,d) in enumerate(feats):
    col=i%3; row=i//3
    x=Emu(int(x0)+col*(int(cw)+int(gx))); y=Emu(int(y0)+row*(int(ch)+int(gy)))
    card=rect(s,x,y,cw,ch,fill=C['cosmos2'],line=C['glow'],line_w=0.6,radius=0.10)
    set_alpha(card,50); set_alpha_line(card,65)
    dot=rect(s,Emu(int(x)+Emu(Inches(0.3))),Emu(int(y)+Emu(Inches(0.32))),Inches(0.16),Inches(0.16),fill=C['glow_l'],shape=MSO_SHAPE.OVAL)
    text(s,Emu(int(x)+Emu(Inches(0.28))),Emu(int(y)+Emu(Inches(0.66))),Emu(int(cw)-Emu(Inches(0.56))),Inches(0.8),
         [[(t,dict(size=14.5,color=C['white'],font=F_HEAD,bold=True))],
          [(d,dict(size=10.5,color=C['mute_d'],font=F_BODY))]],line_spacing=1.0,space_after=2)
home_button(s); page_num(s,7)
transition(s, 'zoom', spd='med')

# =====================================================================
# SLIDE 8 — INSIGHTS & ANALYTICS (overview dashboard)
# =====================================================================
s = slide('bg_dark2.png')
eyebrow(s, Inches(0.85), Inches(0.6), '06 · Insights')
text(s, Inches(0.8), Inches(0.95), Inches(11.6), Inches(0.8),
     [[('Your store’s conversations, made clear', dict(size=32, color=C['white'], font=F_HEAD, bold=True))]])
text(s, Inches(0.8), Inches(1.66), Inches(11.6), Inches(0.5),
     [[('A simple dashboard, built right into your store — no spreadsheets, no guesswork.', dict(size=13.5, color=C['mute_d'], font=F_BODY))]])
panel=rect(s, Inches(0.8), Inches(2.25), Inches(11.73), Inches(4.45), fill=C['cosmos2'], line=C['glow'], line_w=1.0, radius=0.04, shadow=True)
set_alpha(panel,25); set_alpha_line(panel,50)
# KPI tiles
kpis=[('1,240','Conversations this month'),('94%','Answered instantly'),
      ('310','After-hours chats handled'),('4.7 / 5','Average customer rating')]
tw=Inches(2.66); th=Inches(1.15); tx0=Inches(1.05); ty=Inches(2.55); step=Inches(2.84)
for i,(big,lab) in enumerate(kpis):
    x=Emu(int(tx0)+i*int(step))
    tile=rect(s,x,ty,tw,th,fill=C['cosmos'],line=C['glow'],line_w=0.6,radius=0.12)
    set_alpha(tile,35); set_alpha_line(tile,60)
    text(s,Emu(int(x)+Emu(Inches(0.25))),Emu(int(ty)+Emu(Inches(0.16))),Emu(int(tw)-Emu(Inches(0.4))),Inches(0.55),
         [[(big,dict(size=30,color=C['glow_l'],font=F_HEAD,bold=True))]])
    text(s,Emu(int(x)+Emu(Inches(0.25))),Emu(int(ty)+Emu(Inches(0.74))),Emu(int(tw)-Emu(Inches(0.4))),Inches(0.4),
         [[(lab,dict(size=10,color=C['mute_d'],font=F_BODY))]])
# bar chart (left)
text(s,Inches(1.05),Inches(4.05),Inches(5),Inches(0.4),
     [[('Conversations per day',dict(size=13,color=C['white'],font=F_HEAD,bold=True))]])
days=['Mon','Tue','Wed','Thu','Fri','Sat','Sun']; hts=[0.55,0.70,0.62,0.85,1.0,0.92,0.68]
base=Inches(6.2); maxh=Inches(1.45); bw=Inches(0.5); bx0=Inches(1.2); bstep=Inches(0.84)
for i,(d,h) in enumerate(zip(days,hts)):
    bh=Emu(int(int(maxh)*h)); x=Emu(int(bx0)+i*int(bstep))
    col=C['glow_l'] if h==1.0 else C['glow']
    bar=rect(s,x,Emu(int(base)-int(bh)),bw,bh,fill=col,radius=0.2)
    set_alpha(bar, 0 if h==1.0 else 25)
    text(s,Emu(int(x)-Emu(Inches(0.1))),Emu(int(base)+Emu(Inches(0.04))),Inches(0.7),Inches(0.3),
         [[(d,dict(size=8.5,color=C['mute_d'],font=F_MONO))]],align=PP_ALIGN.CENTER)
# handled breakdown (right)
text(s,Inches(7.75),Inches(4.05),Inches(4.5),Inches(0.4),
     [[('How messages were handled',dict(size=13,color=C['white'],font=F_HEAD,bold=True))]])
hb=[('Answered instantly',82,C['glow']),('Passed to your team',18,C['glow_l'])]
trk_w=Inches(4.0)
for i,(lab,pct,col) in enumerate(hb):
    yy=Emu(int(Inches(4.62))+i*int(Inches(0.92)))
    text(s,Inches(7.75),yy,Inches(3.2),Inches(0.3),[[(lab,dict(size=11,color=C['white'],font=F_BODY))]])
    text(s,Inches(10.55),yy,Inches(1.7),Inches(0.3),[[(f'{pct}%',dict(size=12,color=col,font=F_MONO,bold=True))]],align=PP_ALIGN.RIGHT)
    trky=Emu(int(yy)+Emu(Inches(0.34)))
    trk=rect(s,Inches(7.75),trky,trk_w,Inches(0.2),fill=C['cosmos'],radius=0.5); set_alpha(trk,55)
    fillw=Emu(int(int(trk_w)*pct/100))
    fl=rect(s,Inches(7.75),trky,fillw,Inches(0.2),fill=col,radius=0.5)
text(s,Inches(7.75),Inches(6.18),Inches(4.6),Inches(0.4),
     [[('Most asked:  ',dict(size=11,color=C['mute_d'],font=F_BODY)),
       ('“Where’s my order?”',dict(size=11,color=C['glow_l'],font=F_BODY,bold=True))]])
home_button(s); page_num(s,8)
transition(s, 'cover', dir='l')

# =====================================================================
# SLIDE 9 — INSIGHTS IN ACTION
# =====================================================================
s = slide('bg_light.png')
eyebrow(s, Inches(0.85), Inches(0.7), '06 · Insights', color=C['glow'])
text(s, Inches(0.8), Inches(1.05), Inches(11.5), Inches(0.8),
     [[('Turn insight into action', dict(size=34, color=C['ink'], font=F_HEAD, bold=True))]])
text(s, Inches(0.8), Inches(1.78), Inches(11.5), Inches(0.5),
     [[('Spot patterns, staff smarter, and see exactly what turns chats into sales.', dict(size=13.5, color=C['mute_l'], font=F_BODY))]])
icw=Inches(3.85); igap=Inches(0.34); ix0=Inches(0.8); iy=Inches(2.5); ich=Inches(4.0)
cxs=[Emu(int(ix0)+i*(int(icw)+int(igap))) for i in range(3)]
for x in cxs:
    card=rect(s,x,iy,icw,ich,fill=C['white'],line=C['glow'],line_w=1.0,radius=0.08,shadow=True)
    set_alpha_line(card,80)
# Card 1 — heatmap
x=cxs[0]
text(s,Emu(int(x)+Emu(Inches(0.32))),Emu(int(iy)+Emu(Inches(0.28))),Inches(3.3),Inches(0.5),
     [[('Your busiest hours',dict(size=16,color=C['ink'],font=F_HEAD,bold=True))]])
inten=[0.20,0.30,0.25,0.40,0.55,0.45, 0.30,0.45,0.55,0.70,0.85,0.60,
       0.40,0.60,0.75,0.90,1.00,0.80, 0.25,0.40,0.50,0.65,0.55,0.35]
hcols=6; hrows=4; cell=Inches(0.42); cg=Inches(0.1)
hx=Emu(int(x)+Emu(Inches(0.34))); hy=Emu(int(iy)+Emu(Inches(1.0)))
for r in range(hrows):
    for c in range(hcols):
        idx=r*hcols+c
        cxp=Emu(int(hx)+c*(int(cell)+int(cg))); cyp=Emu(int(hy)+r*(int(cell)+int(cg)))
        cl=rect(s,cxp,cyp,cell,cell,fill=C['glow'],radius=0.22)
        set_alpha(cl, 100-int(inten[idx]*100))
text(s,Emu(int(x)+Emu(Inches(0.32))),Emu(int(iy)+Emu(Inches(3.35))),Emu(int(icw)-Emu(Inches(0.6))),Inches(0.6),
     [[('Evenings are your peak — staff up when it matters most.',dict(size=11,color=C['mute_l'],font=F_BODY))]],line_spacing=1.05)
# Card 2 — top questions
x=cxs[1]
text(s,Emu(int(x)+Emu(Inches(0.32))),Emu(int(iy)+Emu(Inches(0.28))),Inches(3.3),Inches(0.5),
     [[('Top questions',dict(size=16,color=C['ink'],font=F_HEAD,bold=True))]])
qs=[('Where’s my order?',34),('Is this in stock?',26),('Shipping & delivery',21),('Returns & exchanges',12)]
qtrk=Inches(3.05)
for i,(q,p) in enumerate(qs):
    yy=Emu(int(iy)+Emu(Inches(1.05))+i*int(Inches(0.72)))
    text(s,Emu(int(x)+Emu(Inches(0.34))),yy,Inches(2.5),Inches(0.3),[[(q,dict(size=11,color=C['ink'],font=F_BODY))]])
    text(s,Emu(int(x)+Emu(Inches(2.7))),yy,Inches(0.85),Inches(0.3),[[(f'{p}%',dict(size=11,color=C['glow'],font=F_MONO,bold=True))]],align=PP_ALIGN.RIGHT)
    trky=Emu(int(yy)+Emu(Inches(0.3)))
    bg=rect(s,Emu(int(x)+Emu(Inches(0.34))),trky,qtrk,Inches(0.14),fill=C['glow'],radius=0.5); set_alpha(bg,88)
    fw=Emu(int(int(qtrk)*p/100))
    fb=rect(s,Emu(int(x)+Emu(Inches(0.34))),trky,fw,Inches(0.14),fill=C['glow'],radius=0.5)
# Card 3 — conversion
x=cxs[2]
text(s,Emu(int(x)+Emu(Inches(0.32))),Emu(int(iy)+Emu(Inches(0.28))),Inches(3.3),Inches(0.5),
     [[('Chats that became sales',dict(size=16,color=C['ink'],font=F_HEAD,bold=True))]])
text(s,Emu(int(x)+Emu(Inches(0.3))),Emu(int(iy)+Emu(Inches(0.95))),Inches(3),Inches(0.9),
     [[('19%',dict(size=52,color=C['glow'],font=F_HEAD,bold=True))]])
text(s,Emu(int(x)+Emu(Inches(0.34))),Emu(int(iy)+Emu(Inches(2.0))),Emu(int(icw)-Emu(Inches(0.6))),Inches(0.6),
     [[('of conversations led to a purchase.',dict(size=11.5,color=C['mute_l'],font=F_BODY))]],line_spacing=1.05)
cbase=Emu(int(iy)+Emu(Inches(3.5))); cmax=Inches(1.0); cbw=Inches(0.55); cbx=Emu(int(x)+Emu(Inches(0.4))); cbs=Inches(0.78)
chs=[0.45,0.65,0.85,1.0]; lbls=['Wk1','Wk2','Wk3','Wk4']
for i,(h,lb) in enumerate(zip(chs,lbls)):
    bh=Emu(int(int(cmax)*h)); bx=Emu(int(cbx)+i*int(cbs))
    rect(s,bx,Emu(int(cbase)-int(bh)),cbw,bh,fill=C['glow_l'],radius=0.2)
    text(s,Emu(int(bx)-Emu(Inches(0.08))),Emu(int(cbase)+Emu(Inches(0.04))),Inches(0.7),Inches(0.3),
         [[(lb,dict(size=8.5,color=C['mute_l'],font=F_MONO))]],align=PP_ALIGN.CENTER)
home_button(s); page_num(s,9,dark=False)
transition(s, 'split', orient='horz', dir='out')

# =====================================================================
# SLIDE 10 — WHY LUMA
# =====================================================================
s = slide('bg_dark2.png')
eyebrow(s, Inches(0.85), Inches(0.85), '07 · Why LUMA')
text(s, Inches(0.8), Inches(1.25), Inches(11.5), Inches(1.4),
     [[('Not a generic chatbot. ', dict(size=38, color=C['white'], font=F_HEAD, bold=True)),
       ('A real teammate.', dict(size=38, color=C['glow_l'], font=F_HEAD, bold=True))],
      [('Built from the ground up for merchants who know that fast answers drive sales.', dict(size=15, color=C['mute_d'], font=F_BODY))]],
     line_spacing=1.05, space_after=8)
# comparison header row
hy=Inches(2.95)
ch_generic=rect(s,Inches(0.8),hy,Inches(5.85),Inches(0.6),fill=C['cosmos'],radius=0.18)
set_alpha(ch_generic,30)
shape_text(ch_generic,[[('GENERIC CHATBOT',dict(size=13,color=C['mute_d'],font=F_MONO,bold=True,spacing=2))]],align=PP_ALIGN.CENTER)
ch_luma=rect(s,Inches(6.95),hy,Inches(5.55),Inches(0.6),fill=C['glow'],radius=0.18,shadow=True)
shape_text(ch_luma,[[('LUMA',dict(size=13,color=C['white'],font=F_MONO,bold=True,spacing=2))]],align=PP_ALIGN.CENTER)
rows=[('Translates word-for-word','Speaks Arabic and Franco-Arabic naturally'),
      ('Guesses order status','Gives accurate, up-to-date order updates'),
      ('Yet another app to open','Works right inside your store'),
      ('Replies sound robotic','Sounds like your brand'),
      ('Gets stuck when confused','Hands off to your team smoothly')]
ry=Emu(int(hy)+Emu(Inches(0.78)))
for g,l in rows:
    gc=rect(s,Inches(0.8),ry,Inches(5.85),Inches(0.6),fill=C['cosmos2'],line=C['mute_d'],line_w=0.4,radius=0.14)
    set_alpha(gc,55); set_alpha_line(gc,80)
    shape_text(gc,[[('✕  ',dict(size=12,color=RGBColor(0xff,0x7a,0x7a),font=F_BODY,bold=True)),(g,dict(size=12.5,color=C['mute_d'],font=F_BODY))]],align=PP_ALIGN.LEFT)
    lc=rect(s,Inches(6.95),ry,Inches(5.55),Inches(0.6),fill=C['cosmos2'],line=C['glow'],line_w=0.8,radius=0.14)
    set_alpha(lc,30); set_alpha_line(lc,40)
    shape_text(lc,[[('✓  ',dict(size=12,color=C['green'],font=F_BODY,bold=True)),(l,dict(size=12.5,color=C['white'],font=F_BODY,bold=True))]],align=PP_ALIGN.LEFT)
    ry=Emu(int(ry)+Emu(Inches(0.7)))
home_button(s); page_num(s,10)
transition(s, 'wipe', dir='u')

# =====================================================================
# SLIDE 9 — PERFORMANCE STATS
# =====================================================================
s = slide('bg_cover.png')
eyebrow(s, Inches(0.85), Inches(0.95), '08 · Proven performance', color=C['glow_xl'])
text(s, Inches(0.8), Inches(1.35), Inches(11.5), Inches(1),
     [[('Numbers that turn conversations into ', dict(size=34, color=C['white'], font=F_HEAD, bold=True)),
       ('revenue', dict(size=34, color=C['glow_l'], font=F_HEAD, bold=True))]])
stats=[('Seconds','Reply time','Customers get an answer almost the moment they ask.'),
       ('3','Languages','Arabic, English and Franco-Arabic, understood automatically.'),
       ('24/7','Always on','Never sleeps — your customers always hear back.'),
       ('3','Channels','Instagram and website chat today, WhatsApp soon.')]
cw=Inches(2.85); gap=Inches(0.3); x0=Inches(0.8); y=Inches(3.0); ch=Inches(3.0)
for i,(big,t,d) in enumerate(stats):
    x=Emu(int(x0)+i*(int(cw)+int(gap)))
    card=rect(s,x,y,cw,ch,fill=C['cosmos2'],line=C['glow'],line_w=1.0,radius=0.08,shadow=True)
    set_alpha(card,30); set_alpha_line(card,45)
    text(s,Emu(int(x)),Emu(int(y)+Emu(Inches(0.45))),cw,Inches(1.2),
         [[(big,dict(size=58,color=C['glow_l'],font=F_HEAD,bold=True))]],align=PP_ALIGN.CENTER)
    text(s,Emu(int(x)+Emu(Inches(0.25))),Emu(int(y)+Emu(Inches(1.75))),Emu(int(cw)-Emu(Inches(0.5))),Inches(1.1),
         [[(t,dict(size=18,color=C['white'],font=F_HEAD,bold=True))],
          [(d,dict(size=11.5,color=C['mute_d'],font=F_BODY))]],align=PP_ALIGN.CENTER,line_spacing=1.05,space_after=4)
home_button(s); page_num(s,11)
transition(s, 'cover', dir='u')

# =====================================================================
# SLIDE 10 — PRICING
# =====================================================================
s = slide('bg_light.png')
eyebrow(s, Inches(0.85), Inches(0.7), '09 · Pricing', color=C['glow'])
text(s, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.8),
     [[('Plans that scale with your store', dict(size=34, color=C['ink'], font=F_HEAD, bold=True))]])
text(s, Inches(0.8), Inches(1.78), Inches(11.5), Inches(0.5),
     [[('Billed through Shopify · 7-day free trial on paid plans · card never touches LUMA', dict(size=13, color=C['mute_l'], font=F_BODY))]])
plans=[('Free','$0','/mo','20 chats / day',['Instagram + Shopify chat','Order tracking','Smart pre-filtering','20 knowledge entries'],False),
       ('Starter','$29','/mo','200 chats / day',['Everything in Free','Custom auto-reply templates','100 knowledge entries','Advanced analytics'],True),
       ('Growth','$79','/mo','600 chats / day',['Everything in Starter','Unlimited knowledge','Priority support','Custom bot personality'],False),
       ('Pro','$149','/mo','1,500 chats / day',['Everything in Growth','Hands-on onboarding','Direct Slack support','Early access features'],False)]
cw=Inches(2.92); gap=Inches(0.18); x0=Inches(0.72); y=Inches(2.55); ch=Inches(4.25)
for name,price,per,limit,fl,pop in plans:
    x=Emu(int(x0)+plans.index([p for p in plans if p[0]==name][0])*(int(cw)+int(gap)))
    idx=[p[0] for p in plans].index(name)
    x=Emu(int(x0)+idx*(int(cw)+int(gap)))
    if pop:
        card=rect(s,x,Emu(int(y)-Emu(Inches(0.15))),cw,Emu(int(ch)+Emu(Inches(0.3))),fill=C['glow'],radius=0.06,shadow=True)
        nm_col=C['white']; pr_col=C['white']; lim_col=C['glow_xl']; ft_col=C['white']; line_chip=True
        tag=rect(s,Emu(int(x)+Emu(Inches(0.85))),Emu(int(y)-Emu(Inches(0.42))),Inches(1.2),Inches(0.4),fill=C['ink'],radius=0.5,shadow=True)
        shape_text(tag,[[('POPULAR',dict(size=9,color=C['glow_xl'],font=F_MONO,bold=True,spacing=1.5))]],align=PP_ALIGN.CENTER)
    else:
        card=rect(s,x,y,cw,ch,fill=C['white'],line=C['glow'],line_w=0.8,radius=0.06,shadow=True)
        set_alpha_line(card,82)
        nm_col=C['ink']; pr_col=C['ink']; lim_col=C['mute_l']; ft_col=C['mute_l']; line_chip=False
    yy=Emu(int(y)+Emu(Inches(0.0 if not pop else 0.12)))
    text(s,Emu(int(x)+Emu(Inches(0.3))),Emu(int(yy)+Emu(Inches(0.25))),Emu(int(cw)-Emu(Inches(0.6))),Inches(0.5),
         [[(name,dict(size=17,color=nm_col,font=F_HEAD,bold=True))]])
    text(s,Emu(int(x)+Emu(Inches(0.3))),Emu(int(yy)+Emu(Inches(0.72))),Emu(int(cw)-Emu(Inches(0.6))),Inches(0.7),
         [[(price,dict(size=38,color=pr_col,font=F_HEAD,bold=True)),(per,dict(size=14,color=lim_col,font=F_BODY))]])
    text(s,Emu(int(x)+Emu(Inches(0.3))),Emu(int(yy)+Emu(Inches(1.5))),Emu(int(cw)-Emu(Inches(0.6))),Inches(0.4),
         [[(limit,dict(size=12,color=lim_col,font=F_MONO,bold=True))]])
    fy=Emu(int(yy)+Emu(Inches(2.0)))
    for f in fl:
        text(s,Emu(int(x)+Emu(Inches(0.3))),fy,Emu(int(cw)-Emu(Inches(0.55))),Inches(0.4),
             [[('✓  ',dict(size=11,color=(C['glow_xl'] if pop else C['glow']),font=F_BODY,bold=True)),(f,dict(size=11,color=ft_col,font=F_BODY))]])
        fy=Emu(int(fy)+Emu(Inches(0.42)))
home_button(s); page_num(s,12,dark=False)
transition(s, 'push', dir='u')

# =====================================================================
# SLIDE 11 — SECURITY & TRUST
# =====================================================================
s = slide('bg_dark.png')
eyebrow(s, Inches(0.85), Inches(0.95), '10 · Security & trust')
text(s, Inches(0.8), Inches(1.35), Inches(11.5), Inches(1),
     [[('Your data is ', dict(size=38, color=C['white'], font=F_HEAD, bold=True)),
       ('protected by design', dict(size=38, color=C['glow_l'], font=F_HEAD, bold=True))]])
sec=[('Your data stays yours','Every store’s information is kept completely separate and private.'),
     ('Safe and secure','Conversations are protected and access is tightly controlled.'),
     ('Payments stay with your store','Billing runs through your store’s trusted, secure checkout.'),
     ('Privacy respected','Customer data-removal requests are handled automatically.')]
cw=Inches(5.85); ch=Inches(1.85); gx=Inches(0.3); gy=Inches(0.3); x0=Inches(0.8); y0=Inches(2.85)
for i,(t,d) in enumerate(sec):
    col=i%2; row=i//2
    x=Emu(int(x0)+col*(int(cw)+int(gx))); y=Emu(int(y0)+row*(int(ch)+int(gy)))
    card=rect(s,x,y,cw,ch,fill=C['cosmos2'],line=C['glow'],line_w=0.75,radius=0.09)
    set_alpha(card,45); set_alpha_line(card,55)
    dot=rect(s,Emu(int(x)+Emu(Inches(0.4))),Emu(int(y)+Emu(Inches(0.42))),Inches(0.2),Inches(0.2),fill=C['glow_l'],shape=MSO_SHAPE.OVAL)
    text(s,Emu(int(x)+Emu(Inches(1.0))),Emu(int(y)+Emu(Inches(0.32))),Emu(int(cw)-Emu(Inches(1.3))),Inches(1.3),
         [[(t,dict(size=17,color=C['white'],font=F_HEAD,bold=True))],
          [(d,dict(size=12.5,color=C['mute_d'],font=F_BODY))]],line_spacing=1.08,space_after=4)
home_button(s); page_num(s,13)
transition(s, 'split', orient='horz', dir='in')

# =====================================================================
# SLIDE 12 — NEXT STEPS / CTA
# =====================================================================
s = slide('bg_cta.png')
s.shapes.add_picture('logo.png', Inches(6.07), Inches(0.95), height=Inches(1.1))
text(s, Inches(1), Inches(2.25), Inches(11.33), Inches(1.5),
     [[('Never lose a sale to silence again', dict(size=42, color=C['white'], font=F_HEAD, bold=True))],
      [('Start free today — live in your store this afternoon.', dict(size=20, color=C['white'], font=F_SERIF, italic=True))]],
     align=PP_ALIGN.CENTER, line_spacing=1.1, space_after=8)
# CTA buttons
b1=rect(s,Inches(3.85),Inches(4.25),Inches(2.7),Inches(0.72),fill=C['white'],radius=0.5,shadow=True)
shape_text(b1,[[('Request access  →',dict(size=15,color=C['glow'],font=F_HEAD,bold=True))]],align=PP_ALIGN.CENTER)
link_url(b1,'https://luma-bot.com')
b2=rect(s,Inches(6.78),Inches(4.25),Inches(2.7),Inches(0.72),fill=None,line=C['white'],line_w=1.5,radius=0.5)
shape_text(b2,[[('Email us',dict(size=15,color=C['white'],font=F_HEAD,bold=True))]],align=PP_ALIGN.CENTER)
link_url(b2,'mailto:hello@luma-bot.com')
# contact line
text(s, Inches(1), Inches(5.5), Inches(11.33), Inches(0.6),
     [[('luma-bot.com', dict(size=15, color=C['white'], font=F_MONO, bold=True)),
       ('     ·     ', dict(size=15, color=C['glow_xl'], font=F_MONO)),
       ('hello@luma-bot.com', dict(size=15, color=C['white'], font=F_MONO, bold=True))]],
     align=PP_ALIGN.CENTER)
text(s, Inches(1), Inches(6.2), Inches(11.33), Inches(0.5),
     [[('Now welcoming a limited number of local businesses', dict(size=12, color=C['glow_xl'], font=F_MONO, spacing=1.5))]],
     align=PP_ALIGN.CENTER)
# back to menu
back=rect(s,Inches(5.86),Inches(6.75),Inches(1.6),Inches(0.45),fill=None,line=C['white'],line_w=1.0,radius=0.5)
set_alpha_line(back,30)
shape_text(back,[[('Back to menu',dict(size=10,color=C['white'],font=F_MONO,bold=True))]],align=PP_ALIGN.CENTER)
link_to_slide(back,1)
transition(s, 'fade', spd='slow')

_resolve_links()
prs.save('proposal/LUMA_Business_Proposal.pptx')
print('SAVED proposal/LUMA_Business_Proposal.pptx with', len(prs.slides._sldIdLst), 'slides')
